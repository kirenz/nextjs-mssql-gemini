# File: backend/app/services/report_service.py

from __future__ import annotations

import calendar
import logging
import asyncio
from typing import Any, Dict, List, Optional, Tuple
import base64
from io import BytesIO

import numpy as np
import pandas as pd
import altair as alt
from vl_convert import vegalite_to_png
from sqlalchemy import text
from statsmodels.tsa.statespace.sarimax import SARIMAX
from pptx import Presentation
from pptx.util import Inches, Pt

from ..database.connection import DatabaseConnection
from ..ai_providers.gemini_provider import GeminiProvider

logger = logging.getLogger(__name__)


class ReportService:
    """Business reporting utilities used by the Reports dashboard."""

    MIN_DATA_POINTS = 24

    def __init__(
        self,
        db: DatabaseConnection,
        ai_provider: Optional[GeminiProvider] = None,
    ) -> None:
        self.db = db
        self.ai_provider = ai_provider

    # ------------------------------------------------------------------
    # Filter helpers
    # ------------------------------------------------------------------
    def get_filter_options(
        self,
        *,
        sales_org: Optional[str] = None,
        country: Optional[str] = None,
        region: Optional[str] = None,
        state: Optional[str] = None,
        product_line: Optional[str] = None,
    ) -> Dict[str, List[str]]:
        """
        Return dropdown options with dependencies similar to the Gradio app.
        """
        try:
            with self.db.engine.connect() as connection:
                sales_orgs = self._with_all(
                    self._fetch_unique(connection, "Sales Organisation")
                )

                countries = self._with_all(
                    self._fetch_unique(
                        connection,
                        "Sales Country",
                        filters={"Sales Organisation": sales_org},
                    )
                )

                regions = self._with_all(
                    self._fetch_unique(
                        connection,
                        "Sales Region",
                        filters={
                            "Sales Organisation": sales_org,
                            "Sales Country": country,
                        },
                    )
                )

                states = self._with_all(
                    self._fetch_unique(
                        connection,
                        "Sales State",
                        filters={
                            "Sales Organisation": sales_org,
                            "Sales Country": country,
                            "Sales Region": region,
                        },
                    )
                )

                cities = self._with_all(
                    self._fetch_unique(
                        connection,
                        "Sales City",
                        filters={
                            "Sales Organisation": sales_org,
                            "Sales Country": country,
                            "Sales Region": region,
                            "Sales State": state,
                        },
                        require_min_points=True,
                    )
                )

                product_lines = self._with_all(
                    self._fetch_unique(connection, "Product Line")
                )

                product_categories = self._with_all(
                    self._fetch_unique(
                        connection,
                        "Product Category",
                        filters={"Product Line": product_line},
                    )
                )

            return {
                "sales_organisations": sales_orgs,
                "countries": countries,
                "regions": regions,
                "states": states,
                "cities": cities,
                "product_lines": product_lines,
                "product_categories": product_categories,
            }
        except Exception as exc:
            logger.exception("Failed to load report filters")
            raise exc

    def _fetch_unique(
        self,
        connection,
        column: str,
        *,
        filters: Optional[Dict[str, Optional[str]]] = None,
        require_min_points: bool = False,
    ) -> List[str]:
        clause, params = self._build_filter_clause(filters or {})
        column_expr = f"[{column}]"
        non_null_clause = f" AND {column_expr} IS NOT NULL"

        if require_min_points:
            query = text(
                f"""
                SELECT {column_expr} as value
                FROM DataSet_Monthly_Sales_and_Quota
                WHERE 1=1 {clause} {non_null_clause}
                GROUP BY {column_expr}
                HAVING COUNT(*) >= :min_points
                ORDER BY {column_expr}
                """
            )
            params["min_points"] = self.MIN_DATA_POINTS
        else:
            query = text(
                f"""
                SELECT DISTINCT {column_expr} as value
                FROM DataSet_Monthly_Sales_and_Quota
                WHERE 1=1 {clause} {non_null_clause}
                ORDER BY {column_expr}
                """
            )

        rows = connection.execute(query, params).fetchall()
        values = [row[0] for row in rows]
        return self._clean_values(values)

    def _build_filter_clause(
        self, filters: Dict[str, Optional[str]]
    ) -> Tuple[str, Dict[str, str]]:
        clause_parts: List[str] = []
        params: Dict[str, str] = {}

        for column, raw_value in filters.items():
            value = self._normalize_filter(raw_value)
            if value is None:
                continue
            param_name = self._param_key(column)
            clause_parts.append(f" AND [{column}] = :{param_name}")
            params[param_name] = value

        return "".join(clause_parts), params

    def _param_key(self, column: str) -> str:
        return column.lower().replace(" ", "_")

    def _normalize_filter(self, value: Optional[str]) -> Optional[str]:
        if value is None:
            return None
        if isinstance(value, str):
            cleaned = value.strip()
            if not cleaned or cleaned.lower() == "all":
                return None
            return cleaned
        return value

    def _clean_values(self, raw_values: List[Any]) -> List[str]:
        seen = set()
        cleaned: List[str] = []
        for value in raw_values:
            if value is None:
                continue
            text_value = str(value).strip()
            if not text_value or text_value in seen:
                continue
            seen.add(text_value)
            cleaned.append(text_value)
        return cleaned

    def _with_all(self, values: List[str]) -> List[str]:
        return ["All"] + values if values else ["All"]

    # ------------------------------------------------------------------
    # Forecast helpers
    # ------------------------------------------------------------------
    def generate_forecast(
        self,
        *,
        sales_org: Optional[str] = None,
        country: Optional[str] = None,
        region: Optional[str] = None,
        state: Optional[str] = None,
        city: Optional[str] = None,
        product_line: Optional[str] = None,
        product_category: Optional[str] = None,
        forecast_periods: int = 12,
        confidence_interval: float = 0.95,
    ) -> Dict[str, Any]:
        """
        Run SARIMAX forecasting and build response payload for the UI.
        """
        return self._build_report_payload(
            sales_org=sales_org,
            country=country,
            region=region,
            state=state,
            city=city,
            product_line=product_line,
            product_category=product_category,
            forecast_periods=forecast_periods,
            confidence_interval=confidence_interval,
        )

    def generate_pptx(
        self,
        *,
        sales_org: Optional[str] = None,
        country: Optional[str] = None,
        region: Optional[str] = None,
        state: Optional[str] = None,
        city: Optional[str] = None,
        product_line: Optional[str] = None,
        product_category: Optional[str] = None,
        forecast_periods: int = 12,
        confidence_interval: float = 0.95,
    ) -> Tuple[bytes, str]:
        """
        Build a PowerPoint report using the same data as the JSON payload.
        """
        payload = self._build_report_payload(
            sales_org=sales_org,
            country=country,
            region=region,
            state=state,
            city=city,
            product_line=product_line,
            product_category=product_category,
            forecast_periods=forecast_periods,
            confidence_interval=confidence_interval,
        )
        ppt_bytes = self._build_pptx_document(payload)
        filename = self._build_report_filename(payload["filters"])
        return ppt_bytes, filename

    def _build_report_payload(
        self,
        *,
        sales_org: Optional[str],
        country: Optional[str],
        region: Optional[str],
        state: Optional[str],
        city: Optional[str],
        product_line: Optional[str],
        product_category: Optional[str],
        forecast_periods: int,
        confidence_interval: float,
    ) -> Dict[str, Any]:
        filtered_data = self._get_filtered_data(
            sales_org=sales_org,
            country=country,
            region=region,
            state=state,
            city=city,
            product_line=product_line,
            product_category=product_category,
        )

        ts = self._prepare_time_series(filtered_data)
        if ts.empty or len(ts) < self.MIN_DATA_POINTS:
            raise ValueError(
                f"Need at least {self.MIN_DATA_POINTS} monthly data points for forecasting."
            )

        model = SARIMAX(
            ts,
            order=(1, 1, 1),
            seasonal_order=(1, 1, 1, 12),
            enforce_stationarity=False,
        )
        results = model.fit(disp=False)

        forecast = results.get_forecast(steps=forecast_periods)
        ci = forecast.conf_int(alpha=1 - confidence_interval)
        ci.columns = ["lower", "upper"]

        historical_agg = (
            filtered_data.groupby("Calendar DueDate")
            .agg({"Revenue EUR": "sum", "Sales Amount": "sum"})
            .reset_index()
            .sort_values("Calendar DueDate")
        )

        seasonality = (
            filtered_data.assign(
                Year=filtered_data["Calendar DueDate"].dt.year,
                Month=filtered_data["Calendar DueDate"].dt.month,
            )
            .groupby(["Year", "Month"])["Revenue EUR"]
            .mean()
            .reset_index()
            .sort_values(["Year", "Month"])
        )

        historical_series = [
            {
                "date": row["Calendar DueDate"].date().isoformat(),
                "revenue": float(row["Revenue EUR"]),
                "sales_amount": float(row["Sales Amount"]),
            }
            for _, row in historical_agg.iterrows()
        ]

        forecast_series = [
            {
                "date": idx.date().isoformat(),
                "forecast": float(value),
                "lower": float(ci.loc[idx, "lower"]),
                "upper": float(ci.loc[idx, "upper"]),
            }
            for idx, value in forecast.predicted_mean.items()
        ]

        seasonality_series = [
            {
                "year": int(row["Year"]),
                "month": int(row["Month"]),
                "label": calendar.month_abbr[int(row["Month"])],
                "revenue": float(row["Revenue EUR"]),
            }
            for _, row in seasonality.iterrows()
        ]

        table_rows = [
            {
                "date": entry["date"],
                "forecast": entry["forecast"],
                "lower": entry["lower"],
                "upper": entry["upper"],
            }
            for entry in forecast_series
        ]

        residuals = results.resid
        non_zero_mask = (ts != 0) & (~residuals.isna())
        if non_zero_mask.any():
            mape = float(
                np.mean(
                    np.abs(residuals[non_zero_mask] / ts[non_zero_mask])
                )
                * 100
            )
        else:
            mape = None

        filters_used = {
            "Sales Organization": sales_org or "All",
            "Country": country or "All",
            "Region": region or "All",
            "State": state or "All",
            "City": city or "All",
            "Product Line": product_line or "All",
            "Product Category": product_category or "All",
        }

        summary_lines = [
            "Forecast Analysis Summary",
            "-------------------------",
            "Applied Filters:",
            *[f"- {key}: {value}" for key, value in filters_used.items()],
            "",
            f"Data Points Analyzed: {len(filtered_data)}",
            f"Forecast Periods: {forecast_periods}",
            f"Confidence Interval: {confidence_interval * 100:.0f}%",
            (
                f"MAPE: {mape:.2f}%"
                if mape is not None
                else "MAPE: Not enough data to calculate"
            ),
            f"Latest Historical Value: {ts.iloc[-1]:,.2f} EUR",
            (
                f"Latest Forecast Value: {forecast.predicted_mean.iloc[-1]:,.2f} EUR"
                if not forecast.predicted_mean.empty
                else ""
            ),
        ]

        summary = "\n".join(line for line in summary_lines if line != "")

        charts = self._build_charts(
            historical_df=historical_agg,
            forecast_series=forecast_series,
            historical_series=historical_series,
            seasonality_series=seasonality_series,
        )

        metrics_payload = {
            "data_points": len(filtered_data),
            "forecast_periods": forecast_periods,
            "confidence_interval": confidence_interval,
            "mape": round(mape, 2) if mape is not None else None,
            "latest_historical": float(ts.iloc[-1]),
            "latest_forecast": float(forecast.predicted_mean.iloc[-1])
            if not forecast.predicted_mean.empty
            else None,
        }

        explanation = self._build_ai_explanation(
            summary=summary,
            filters=filters_used,
            metrics=metrics_payload,
            forecast_series=forecast_series,
        )

        return {
            "summary": summary,
            "filters": filters_used,
            "metrics": metrics_payload,
            "historical_series": historical_series,
            "forecast_series": forecast_series,
            "seasonality_series": seasonality_series,
            "forecast_table": table_rows,
            "charts": charts,
            "explanation": explanation,
        }

    def _get_filtered_data(
        self,
        *,
        sales_org: Optional[str] = None,
        country: Optional[str] = None,
        region: Optional[str] = None,
        state: Optional[str] = None,
        city: Optional[str] = None,
        product_line: Optional[str] = None,
        product_category: Optional[str] = None,
    ) -> pd.DataFrame:
        query = """
            SELECT 
                [Calendar DueDate],
                [Revenue EUR],
                [Sales Amount],
                [Sales Organisation],
                [Sales Country],
                [Sales Region],
                [Sales State],
                [Sales City],
                [Product Line],
                [Product Category]
            FROM DataSet_Monthly_Sales_and_Quota
            WHERE 1=1
        """

        clause, params = self._build_filter_clause(
            {
                "Sales Organisation": sales_org,
                "Sales Country": country,
                "Sales Region": region,
                "Sales State": state,
                "Sales City": city,
                "Product Line": product_line,
                "Product Category": product_category,
            }
        )

        query = query + clause + " ORDER BY [Calendar DueDate]"

        with self.db.engine.connect() as connection:
            df = pd.read_sql_query(text(query), connection, params=params)

        if df.empty:
            raise ValueError("No data found for the selected filters.")

        if len(df) < self.MIN_DATA_POINTS:
            raise ValueError(
                f"Insufficient data for forecast. Found only {len(df)} rows but "
                f"at least {self.MIN_DATA_POINTS} are required."
            )

        df["Calendar DueDate"] = pd.to_datetime(df["Calendar DueDate"])
        return df

    def _prepare_time_series(self, df: pd.DataFrame) -> pd.Series:
        series = (
            df.groupby("Calendar DueDate")["Revenue EUR"]
            .sum()
            .sort_index()
        )
        if series.empty:
            return series

        date_range = pd.date_range(
            start=series.index.min(),
            end=series.index.max(),
            freq="M",
        )
        series = series.reindex(date_range, fill_value=0)
        return series

    # ------------------------------------------------------------------
    # Chart helpers
    # ------------------------------------------------------------------
    def _build_charts(
        self,
        *,
        historical_df: pd.DataFrame,
        forecast_series: List[Dict[str, Any]],
        historical_series: List[Dict[str, Any]],
        seasonality_series: List[Dict[str, Any]],
    ) -> Dict[str, Optional[str]]:
        try:
            return {
                "historical": self._chart_to_data_url(
                    self._historical_chart(historical_df)
                ),
                "forecast": self._chart_to_data_url(
                    self._forecast_chart(historical_series, forecast_series)
                ),
                "seasonal": self._chart_to_data_url(
                    self._seasonality_chart(seasonality_series)
                ),
            }
        except Exception as exc:
            logger.warning("Failed to build charts: %s", exc)
            return {"historical": None, "forecast": None, "seasonal": None}

    def _chart_to_data_url(self, chart: alt.Chart) -> str:
        png_bytes = vegalite_to_png(chart.to_dict())
        encoded = base64.b64encode(png_bytes).decode("utf-8")
        return f"data:image/png;base64,{encoded}"

    def _historical_chart(self, historical_df: pd.DataFrame) -> alt.Chart:
        melted = historical_df.rename(
            columns={"Revenue EUR": "Revenue", "Sales Amount": "Sales"}
        ).melt(
            id_vars="Calendar DueDate",
            value_vars=["Revenue", "Sales"],
            var_name="Metric",
            value_name="Value",
        )

        chart = (
            alt.Chart(melted)
            .mark_line()
            .encode(
                x=alt.X("Calendar DueDate:T", title="Date"),
                y=alt.Y("Value:Q", title="Amount (EUR)"),
                color=alt.Color("Metric:N", title=""),
            )
            .properties(width=500, height=300, title="Historical Performance")
        )
        return chart

    def _build_pptx_document(self, payload: Dict[str, Any]) -> bytes:
        prs = Presentation()

        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Sales Forecast Report"
        subtitle = title_slide.placeholders[1]
        subtitle.text = "Generated via Analytics AI Platform"

        # Summary slide
        summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
        summary_slide.shapes.title.text = "Executive Summary"
        summary_body = summary_slide.shapes.placeholders[1].text_frame
        summary_body.text = payload["summary"].strip()

        # KPI slide
        kpi_slide = prs.slides.add_slide(prs.slide_layouts[5])
        kpi_slide.shapes.title.text = "Key Metrics"
        metrics = payload["metrics"]
        text_box = kpi_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
        tf = text_box.text_frame
        tf.word_wrap = True
        lines = [
            f"Data points analyzed: {metrics['data_points']:,}",
            f"Forecast horizon: {metrics['forecast_periods']} months",
            f"Confidence interval: {metrics['confidence_interval'] * 100:.0f}%",
            f"MAPE: {metrics['mape']:.2f}%"
            if metrics["mape"] is not None
            else "MAPE: n/a",
            f"Latest actual: {metrics['latest_historical']:,.2f} EUR",
            f"Latest forecast: {metrics['latest_forecast']:,.2f} EUR"
            if metrics["latest_forecast"] is not None
            else "Latest forecast: n/a",
        ]
        for idx, line in enumerate(lines):
            para = tf.paragraphs[idx] if idx < len(tf.paragraphs) else tf.add_paragraph()
            para.text = line
            para.level = 0
            para.font.size = Pt(16)

        # Chart slides
        chart_titles = [
            ("Historical Performance", payload["charts"].get("historical")),
            ("Forecast with Confidence Bands", payload["charts"].get("forecast")),
            ("Seasonal Patterns", payload["charts"].get("seasonal")),
        ]
        for title, data_url in chart_titles:
            if not data_url:
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title
            image_bytes = self._decode_data_url(data_url)
            slide.shapes.add_picture(
                BytesIO(image_bytes),
                Inches(0.5),
                Inches(1.5),
                width=Inches(9),
            )

        # Gemini insights
        explanation = payload.get("explanation")
        if explanation:
            insights_slide = prs.slides.add_slide(prs.slide_layouts[5])
            insights_slide.shapes.title.text = "Gemini Insights"
            box = insights_slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)
            )
            tf = box.text_frame
            for idx, line in enumerate(explanation.splitlines()):
                para = tf.paragraphs[idx] if idx < len(tf.paragraphs) else tf.add_paragraph()
                para.text = line
                para.level = 0
                para.font.size = Pt(16)

        # Forecast table slide (top rows)
        forecast_rows = payload["forecast_table"][:5]
        if forecast_rows:
            table_slide = prs.slides.add_slide(prs.slide_layouts[5])
            table_slide.shapes.title.text = "Forecast Highlights"
            rows = len(forecast_rows) + 1
            cols = 4
            table = table_slide.shapes.add_table(
                rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(3.5)
            ).table
            headers = ["Date", "Forecast", "Lower", "Upper"]
            for idx, header in enumerate(headers):
                cell = table.cell(0, idx)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
            for row_idx, row in enumerate(forecast_rows, start=1):
                table.cell(row_idx, 0).text = row["date"]
                table.cell(row_idx, 1).text = f"{row['forecast']:,.0f}"
                table.cell(row_idx, 2).text = f"{row['lower']:,.0f}"
                table.cell(row_idx, 3).text = f"{row['upper']:,.0f}"

        output = BytesIO()
        prs.save(output)
        return output.getvalue()

    def _decode_data_url(self, data_url: str) -> bytes:
        if "," in data_url:
            data_url = data_url.split(",", 1)[1]
        return base64.b64decode(data_url)

    def _build_report_filename(self, filters: Dict[str, str]) -> str:
        parts = [value for value in filters.values() if value and value != "All"]
        suffix = "_".join(part.replace(" ", "_") for part in parts) if parts else "global"
        return f"sales_forecast_{suffix}.pptx"

    def _forecast_chart(
        self,
        historical_series: List[Dict[str, Any]],
        forecast_series: List[Dict[str, Any]],
    ) -> alt.Chart:
        hist_df = pd.DataFrame(historical_series)
        hist_df["Type"] = "Historical"
        hist_df.rename(columns={"revenue": "Value"}, inplace=True)
        hist_df = hist_df[["date", "Value", "Type"]]

        forecast_df = pd.DataFrame(forecast_series)
        if forecast_df.empty:
            forecast_df = pd.DataFrame(
                {"date": [], "Value": [], "Type": [], "lower": [], "upper": []}
            )
        else:
            forecast_df = forecast_df.rename(columns={"forecast": "Value"})
            forecast_df["Type"] = "Forecast"

        combined = pd.concat([hist_df, forecast_df], ignore_index=True)
        combined["date"] = pd.to_datetime(combined["date"])

        ci_df = pd.DataFrame(forecast_series)
        if not ci_df.empty:
            ci_df["date"] = pd.to_datetime(ci_df["date"])

        line_chart = (
            alt.Chart(combined)
            .mark_line()
            .encode(
                x=alt.X("date:T", title="Date"),
                y=alt.Y("Value:Q", title="Revenue (EUR)"),
                color=alt.Color("Type:N", title=""),
            )
        )

        if ci_df.empty:
            return line_chart.properties(width=500, height=300, title="Forecast")

        band = (
            alt.Chart(ci_df)
            .mark_area(opacity=0.2)
            .encode(
                x="date:T",
                y="lower:Q",
                y2="upper:Q",
            )
        )

        return (band + line_chart).properties(width=500, height=300, title="Forecast")

    def _seasonality_chart(
        self, seasonality_series: List[Dict[str, Any]]
    ) -> alt.Chart:
        df = pd.DataFrame(seasonality_series)
        if df.empty:
            df = pd.DataFrame({"label": [], "revenue": [], "year": []})
        chart = (
            alt.Chart(df)
            .mark_line()
            .encode(
                x=alt.X("label:N", title="Month"),
                y=alt.Y("revenue:Q", title="Average Revenue (EUR)"),
                color=alt.Color("year:N", title="Year"),
            )
            .properties(width=500, height=300, title="Seasonal Patterns")
        )
        return chart

    def _build_ai_explanation(
        self,
        *,
        summary: str,
        filters: Dict[str, str],
        metrics: Dict[str, Any],
        forecast_series: List[Dict[str, Any]],
    ) -> Optional[str]:
        if not self.ai_provider:
            return None

        sample_forecast = forecast_series[:5]
        prompt = (
            "You are a senior business analyst. "
            "Craft a concise narrative (2 short paragraphs) explaining the forecast "
            "and then list the top 3 KPIs as bullets. "
            "Focus on business implications, drivers, and any risks.\n\n"
            f"Filters: {filters}\n"
            f"Key metrics: {metrics}\n"
            f"Upcoming forecast points: {sample_forecast}\n"
            f"Technical summary:\n{summary}\n"
            "Respond with Markdown."
        )

        try:
            text = asyncio.run(self.ai_provider.generate_analysis(prompt))
            return text.strip()
        except Exception as exc:
            logger.warning("Gemini explanation failed: %s", exc)
            return None
